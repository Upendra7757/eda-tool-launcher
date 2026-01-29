import os
import re
import uuid
import shutil
import subprocess

from django.conf import settings
from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse, HttpResponse, FileResponse
from django.views.decorators.csrf import csrf_exempt
from .services import execute_tool_run
from .models import ToolRun, Tool
from django.utils import timezone

from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import viewsets
from launcher.models import ToolRun
from django.shortcuts import redirect

from .models import Category, Tool, ToolRun, Presentation
from .serializers import ToolSerializer
from django.utils import timezone
from .models import Presentation, Slide
from django.shortcuts import render, get_object_or_404
from django.conf import settings
from pathlib import Path

from .models import Presentation
from launcher.models import SlideItem




# =====================================================
# API VIEWSET
# =====================================================

class ToolViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = Tool.objects.filter(visible=True).order_by("name")
    serializer_class = ToolSerializer


# =====================================================
# FRONTEND PAGES
# =====================================================

def home(request):
    categories = Category.objects.all().order_by("name")
    return render(request, "launcher/home.html", {
        "categories": categories
    })

"""
def category_page(request, slug):
    category = get_object_or_404(Category, slug=slug)
    tools = category.tools.filter(visible=True).order_by("name")
    return render(request, "launcher/category.html", {
        "category": category,
        "tools": tools
    })
"""

def category_page(request, slug):
    category = get_object_or_404(Category, slug=slug)
    tools = category.tools.filter(visible=True).order_by("name")

    return render(request, "launcher/category.html", {
        "category": category,
        "tools": tools,
        "active_category": category.slug,   # ⭐ KEY LINE
    })

def tool_page(request, slug):
    """
    Entry page for any tool.
    Handles WEB redirects internally without breaking UI.
    """
    tool = get_object_or_404(Tool, slug=slug)

    mode = request.GET.get("mode")
    from_tool = request.GET.get("from")

    # -----------------------------
    # VERILATOR WEB WORKSPACE
    # -----------------------------
    if mode == "web" and tool.slug == "verilator":
        return render(request, "launcher/verilator.html", {
            "tool": tool,
            "envs": tool.envs.all(),
            "webmode": True,
            "launched_from": from_tool
        })

    # -----------------------------
    # KLAYOUT WEB WORKSPACE
    # -----------------------------
    if mode == "web" and tool.slug == "klayout":
        return render(request, "launcher/klayout_workspace.html", {
            "tool": tool,
            "webmode": True,
            "launched_from": from_tool
        })

    # -----------------------------
    # DEFAULT TOOL PAGE
    # -----------------------------
    return render(request, "launcher/tool.html", {
        "tool": tool,
        "envs": tool.envs.all()
    })




# =====================================================
# FILE SERVER (UPLOADS)
# =====================================================

def serve_upload(request, path):
    base = os.path.join(settings.BASE_DIR, "uploads")
    full = os.path.abspath(os.path.join(base, path))

    if not full.startswith(os.path.abspath(base)):
        return HttpResponse("Forbidden", status=403)

    if not os.path.exists(full):
        return HttpResponse("Not Found", status=404)

    return FileResponse(open(full, "rb"))


# =====================================================
# VERILATOR HELPERS
# =====================================================

def fix_verilog_module_name(file_path, new_name):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    module_pattern = re.compile(r"(\bmodule\s+)([A-Za-z_][A-Za-z0-9_]*)")
    match = module_pattern.search(text)

    if not match:
        return False

    text = module_pattern.sub(r"\1" + new_name, text, 1)

    if "verilator lint_off DECLFILENAME" not in text:
        text = (
            "/* verilator lint_off DECLFILENAME */\n"
            + text +
            "\n/* verilator lint_on DECLFILENAME */\n"
        )

    if not text.endswith("\n"):
        text += "\n"

    with open(file_path, "w", encoding="utf-8") as f:
        f.write(text)

    return True


def patch_sim_main(path, module_name):
    with open(path, "r") as f:
        txt = f.read()

    txt = txt.replace("VMODULE_NAME", f"V{module_name}")

    with open(path, "w") as f:
        f.write(txt)


# =====================================================
# VERILATOR COMPILE + RUN (WSL)
# =====================================================
"""
@csrf_exempt
def launch_tool(request):
    if request.method != "POST":
        return Response({"ok": False, "error": "POST required"}, status=400)

    tool_id = request.data.get("tool_id")
    upload = request.data.get("upload")

    tool = get_object_or_404(Tool, id=tool_id)

    BASE = str(settings.BASE_DIR)
    win_uploads = os.path.join(BASE, "uploads")
    wsl_uploads = "/mnt/c" + BASE[2:].replace("\\", "/") + "/uploads"

    win_source = os.path.join(win_uploads, upload)
    if not os.path.exists(win_source):
        return Response({"ok": False, "error": "Uploaded file missing"}, status=404)

    folder = uuid.uuid4().hex
    win_workdir = os.path.join(win_uploads, folder)
    wsl_workdir = f"{wsl_uploads}/{folder}"
    os.makedirs(win_workdir, exist_ok=True)

    shutil.copy(win_source, win_workdir)

    module_name = f"top_{os.path.splitext(upload)[0]}"
    fix_verilog_module_name(os.path.join(win_workdir, upload), module_name)

    sim_src = os.path.join(BASE, "launcher", "verilator_assets", "sim_main.cpp")
    sim_dst = os.path.join(win_workdir, "sim_main.cpp")
    shutil.copy(sim_src, sim_dst)
    patch_sim_main(sim_dst, module_name)

    compile_cmd = (
        f"cd {wsl_workdir} && "
        f"verilator -Wall --Wno-EOFNEWLINE --trace --cc {upload} "
        f"--top-module {module_name} "
        f"--exe sim_main.cpp && "
        f"make -C obj_dir -f V{module_name}.mk V{module_name}"
    )

    run_cmd = f"cd {wsl_workdir} && ./obj_dir/V{module_name}"

    def run_bash(cmd):
        p = subprocess.Popen(
            ["wsl", "bash", "-lc", cmd],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )
        out, err = p.communicate()
        return out.decode(), err.decode()

    compile_out, compile_err = run_bash(compile_cmd)
    sim_out, sim_err = run_bash(run_cmd)

    vcd = None
    for f in os.listdir(win_workdir):
        if f.endswith(".vcd"):
            vcd = f"/uploads/{folder}/{f}"

    return Response({
        "ok": True,
        "stdout": compile_out + sim_out,
        "stderr": compile_err + sim_err,
        "vcd": vcd
    })
"""


# =====================================================
# UPLOAD WRAPPER (VERILATOR)
# =====================================================
"""
@api_view(["POST"])
def run_tool(request, slug):
    tool = get_object_or_404(Tool, slug=slug)
    upload = request.FILES.get("file")

    uploads_dir = os.path.join(settings.BASE_DIR, "uploads")
    os.makedirs(uploads_dir, exist_ok=True)

    filename = None
    if upload:
        filename = f"{uuid.uuid4().hex}_{upload.name}"
        with open(os.path.join(uploads_dir, filename), "wb") as f:
            for chunk in upload.chunks():
                f.write(chunk)

    dummy = type("Dummy", (), {})()
    dummy.method = "POST"
    dummy.data = {"tool_id": tool.id, "upload": filename}

    resp = launch_tool(dummy)
    return JsonResponse(resp.data)
"""


# =====================================================
# DESKTOP LAUNCH (WSL SAFE)
# =====================================================


# =====================================================
# WEB LAUNCH (PER TOOL)
# =====================================================
"""
@csrf_exempt
def launch_web(request, slug):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    if slug == "verilator":
        return JsonResponse({
            "ok": True,
            "url": "/tool/verilator/?mode=web&from=launcher"
        })

    if slug == "klayout":
        return JsonResponse({
            "ok": True,
            "url": "/tool/klayout/?mode=web&from=launcher"
        })

    return JsonResponse({
        "ok": False,
        "error": "This tool has no Web Mode"
    })
"""
"""
@csrf_exempt
def klayout_run(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    upload = request.FILES.get("file")
    if not upload:
        return JsonResponse({"ok": False, "error": "No GDS file uploaded"}, status=400)

    base_dir = settings.BASE_DIR
    upload_dir = os.path.join(base_dir, "uploads", "klayout")
    os.makedirs(upload_dir, exist_ok=True)

    filename = f"{uuid.uuid4().hex}_{upload.name}"
    full_path = os.path.join(upload_dir, filename)

    with open(full_path, "wb") as f:
        for chunk in upload.chunks():
            f.write(chunk)

    wsl_path = "/mnt/c" + full_path.replace("C:", "").replace("\\", "/")

    subprocess.Popen(
        ["wsl", "bash", "-lc", f"klayout {wsl_path}"],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL
    )

    return JsonResponse({
        "ok": True,
        "message": "KLayout opened with GDS file"
    })
"""

@csrf_exempt
def launch_web(request, slug):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    # Verilator → Web simulation workspace
    if slug == "verilator":
        return JsonResponse({
            "ok": True,
            "url": "/tool/verilator/?mode=web&from=launcher"
        })

    # KLayout → Web workspace (upload + desktop launch)
    if slug == "klayout":
        return JsonResponse({
            "ok": True,
            "url": "/tool/klayout/?mode=web&from=launcher"
        })

    return JsonResponse({
        "ok": False,
        "error": "This tool does not support Web Mode"
    }, status=400)



@csrf_exempt
def launch_web(request, slug):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    if slug == "klayout":
        return JsonResponse({
            "ok": True,
            "url": "/tool/klayout/?mode=web"
        })

    if slug == "verilator":
        return JsonResponse({
            "ok": True,
            "url": "/tool/verilator/?mode=web"
        })

    return JsonResponse({"ok": False, "error": "No web mode"}, status=400)



@csrf_exempt
def launch_desktop(request, slug):
    if request.method != "POST":
        return JsonResponse(
            {"ok": False, "error": "POST required"},
            status=400
        )

    tool = get_object_or_404(Tool, slug=slug)

    # Prefer Linux executable (WSL GUI tools)
    exe = tool.linux_executable_path or tool.windows_executable_path

    if not exe:
        return JsonResponse(
            {"ok": False, "error": "Executable path not configured"},
            status=400
        )

    # --------------------------------------------------
    # Tool-specific launch flags
    # --------------------------------------------------
    if tool.slug == "klayout":
        # Start KLayout in EDITOR MODE
        cmd = f"{exe} -e"
    else:
        # Default launch
        cmd = exe

    try:
        subprocess.Popen(
            ["wsl", "bash", "-lc", cmd],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            start_new_session=True
        )

        return JsonResponse({
            "ok": True,
            "message": f"{tool.name} launched successfully"
        })

    except Exception as e:
        return JsonResponse(
            {"ok": False, "error": str(e)},
            status=500
        )
    
"""
from django.utils import timezone
from .models import ToolRun

@csrf_exempt
def launch_tool(request):

    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    tool = get_object_or_404(Tool, slug="verilator")

    try:
        # Run actual simulation here
        stdout = "simulation output"
        stderr = ""

        run = execute_tool_run(
            tool=tool,
            user=request.user if request.user.is_authenticated else None,
            input_file=request.FILES.get("file").name if request.FILES else None,
            stdout=stdout,
            stderr=stderr,
            status="success"
        )

        return JsonResponse({
            "ok": True,
            "run_id": run.id,
            "stdout": stdout,
            "stderr": stderr
        })

    except Exception as e:
        run = execute_tool_run(
            tool=tool,
            user=request.user if request.user.is_authenticated else None,
            status="failed",
            stderr=str(e)
        )
        return JsonResponse({"ok": False, "error": str(e)}, status=500)

"""
from .models import ToolRun

"""
@csrf_exempt
def klayout_run(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    upload = request.FILES.get("file")
    if not upload:
        return JsonResponse({"ok": False, "error": "No GDS uploaded"}, status=400)

    tool = get_object_or_404(Tool, slug="klayout")

    run = ToolRun.objects.create(
        tool=tool,
        user=request.user if request.user.is_authenticated else None,
        input_file=upload.name,
        status="running"
    )

    try:
        base = settings.BASE_DIR
        upload_dir = os.path.join(base, "uploads", "klayout")
        os.makedirs(upload_dir, exist_ok=True)

        filename = f"{uuid.uuid4().hex}_{upload.name}"
        full_path = os.path.join(upload_dir, filename)

        with open(full_path, "wb") as f:
            for chunk in upload.chunks():
                f.write(chunk)

        run.status = "success"
        run.completed_at = timezone.now()
        run.save()

        return JsonResponse({
            "ok": True,
            "run_id": run.id,
            "message": "GDS uploaded successfully",
            "path": full_path
        })

    except Exception as e:
        run.status = "failed"
        run.stderr = str(e)
        run.completed_at = timezone.now()
        run.save()

        return JsonResponse({"ok": False, "error": str(e)}, status=500)
"""
"""
@csrf_exempt
def klayout_run(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    upload = request.FILES.get("file")
    if not upload:
        return JsonResponse({"ok": False, "error": "No GDS uploaded"}, status=400)

    upload_dir = os.path.join(settings.BASE_DIR, "uploads", "klayout")
    os.makedirs(upload_dir, exist_ok=True)

    filename = f"{uuid.uuid4().hex}_{upload.name}"
    full_path = os.path.join(upload_dir, filename)

    with open(full_path, "wb") as f:
        for chunk in upload.chunks():
            f.write(chunk)

    # Convert Windows path → WSL path
    wsl_path = "/mnt/c" + full_path.replace("C:", "").replace("\\", "/")

    # IMPORTANT: -e = Editor mode
    subprocess.Popen(
        ["wsl", "bash", "-lc", f"klayout -e '{wsl_path}'"],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL
    )

    return JsonResponse({
        "ok": True,
        "message": "GDS opened in KLayout editor",
    })
"""
def view_run_logs(request, run_id):
    run = get_object_or_404(ToolRun, id=run_id)

    return render(
        request,
        "launcher/run_detail.html",   # ✅ matches your folder
        {
            "run": run
        }
    )
"""
from django.utils import timezone
from .models import ToolRun, LayoutMetadata

@csrf_exempt
def klayout_run(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    upload = request.FILES.get("file")
    if not upload:
        return JsonResponse({"ok": False, "error": "No GDS uploaded"}, status=400)

    run_dir = os.path.join(settings.BASE_DIR, "runs", uuid.uuid4().hex)
    os.makedirs(run_dir, exist_ok=True)

    gds_path = os.path.join(run_dir, upload.name)
    with open(gds_path, "wb") as f:
        for chunk in upload.chunks():
            f.write(chunk)

    run = ToolRun.objects.create(
        tool=Tool.objects.get(slug="klayout"),
        user=request.user if request.user.is_authenticated else None,
        input_file=upload.name,
        run_dir=run_dir,
        status="running",
    )

    wsl_gds = "/mnt/c" + gds_path.replace("C:", "").replace("\\", "/")
    wsl_run = "/mnt/c" + run_dir.replace("C:", "").replace("\\", "/")

    subprocess.Popen([
        "wsl", "bash", "-lc",
        f"klayout -e '{wsl_gds}' && "
        f"klayout -b -r scripts/klayout_extract.py '{wsl_gds}' '{wsl_run}'"
    ])

    return JsonResponse({
        "ok": True,
        "run_id": run.id,
        "message": "Opened in KLayout + metadata generation started"
    })
    run.status = "success"
    run.completed_at = timezone.now()
    run.save()

    # 🔥 AUTO-GENERATE PNG
    generate_klayout_png(run, full_path)
    # 🔥 AUTO-CREATE PRESENTATION (already added)
    create_presentation_for_run(run)

    return run

import os
import subprocess
from django.conf import settings
from .models import LayoutMetadata

def generate_klayout_png(run, gds_path):
    out_dir = os.path.join(settings.BASE_DIR, "runs", str(run.id))
    os.makedirs(out_dir, exist_ok=True)

    png_path = os.path.join(out_dir, "layout.png")

    script = os.path.join(settings.BASE_DIR, "scripts", "klayout_export_png.py")

    subprocess.run([
        "wsl", "bash", "-lc",
        f"klayout -b -r {script} {gds_path} {png_path}"
    ], check=True)

    metadata, _ = LayoutMetadata.objects.get_or_create(run=run)
    metadata.png_preview = png_path
    metadata.save()
"""
import os
import uuid
import time
import subprocess

from django.conf import settings
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.shortcuts import render, get_object_or_404

from .models import ToolRun, Tool, RunArtifact


# ---------------------------
# Helper: Windows → WSL path
# ---------------------------
def wsl_path(p):
    return "/mnt/c" + p.replace("C:", "").replace("\\", "/")
"""
@csrf_exempt
def klayout_run(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    upload = request.FILES.get("file")
    if not upload:
        return JsonResponse({"ok": False, "error": "No GDS uploaded"}, status=400)

    # -------------------------
    # Create run directory
    # -------------------------
    run_dir_name = uuid.uuid4().hex
    run_dir = os.path.join(settings.MEDIA_ROOT, "runs", run_dir_name)

    #run_dir = os.path.join(settings.BASE_DIR, "uploads", "runs", run_dir_name)
    os.makedirs(run_dir, exist_ok=True)

    gds_path = os.path.join(run_dir, "generated_design.gds")
    with open(gds_path, "wb") as f:
        for chunk in upload.chunks():
            f.write(chunk)

    # -------------------------
    # Create ToolRun
    # -------------------------
    run = ToolRun.objects.create(
        tool=Tool.objects.get(slug="klayout"),
        user=request.user if request.user.is_authenticated else None,
        input_file=upload.name,
        run_dir=run_dir,
        status="running",
    )

    # -------------------------
    # Output paths
    # -------------------------
    png_path = os.path.join(run_dir, "preview.png")
    meta_path = os.path.join(run_dir, "metadata.json")
    log_path = os.path.join(run_dir, "klayout.log")

    # -------------------------
    # Run KLayout (NEVER raise)
    # -------------------------
    cmd = (
        f"export KLAYOUT_GDS='{wsl_path(gds_path)}' && "
        f"export KLAYOUT_PNG='{wsl_path(png_path)}' && "
        f"export KLAYOUT_META='{wsl_path(meta_path)}' && "
        f"klayout -b -r scripts/klayout_extract.py "
        f"> '{wsl_path(log_path)}' 2>&1"
    )

    proc =subprocess.run(
        ["wsl", "bash", "-lc", cmd],
        #check=False
        #capture_output=True,
        text=True
    )

    

    
    # -------------------------
    # Register artifacts (ALWAYS)
    # -------------------------
    register_klayout_artifacts(run)

    # -------------------------
    # Create presentation + slides
    # -------------------------
    presentation = auto_create_presentation(run)

    # -------------------------
    # Attach artifacts
    # -------------------------
    auto_attach_artifacts_to_slides(presentation, run)

    return JsonResponse({
        "ok": True,
        "message": "KLayout batch run completed",
        "run_id": str(run.id),
        "presentation_id": presentation.id,
        "redirect": f"/presentation/{presentation.id}/"
    })
"""
@csrf_exempt
def klayout_run(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "POST required"}, status=400)

    upload = request.FILES.get("file")
    if not upload:
        return JsonResponse({"ok": False, "error": "No GDS uploaded"}, status=400)

    # -------------------------
    # Create run directory
    # -------------------------
    run_dir_name = uuid.uuid4().hex
    run_dir = os.path.join(settings.MEDIA_ROOT, "runs", run_dir_name)
    os.makedirs(run_dir, exist_ok=True)

    gds_path = os.path.join(run_dir, "generated_design.gds")
    with open(gds_path, "wb") as f:
        for chunk in upload.chunks():
            f.write(chunk)

    # -------------------------
    # Create ToolRun
    # -------------------------
    run = ToolRun.objects.create(
        tool=Tool.objects.get(slug="klayout"),
        user=request.user if request.user.is_authenticated else None,
        input_file=upload.name,
        run_dir=run_dir,
        status="running",
    )

    # -------------------------
    # Output paths
    # -------------------------
    png_path = os.path.join(run_dir, "preview.png")
    meta_path = os.path.join(run_dir, "metadata.json")
    log_path = os.path.join(run_dir, "klayout.log")

    # -------------------------
    # Run KLayout batch (WSL)
    # -------------------------
    cmd = (
        f"export KLAYOUT_GDS='{wsl_path(gds_path)}' && "
        f"export KLAYOUT_PNG='{wsl_path(png_path)}' && "
        f"export KLAYOUT_META='{wsl_path(meta_path)}' && "
        f"klayout -b -r scripts/klayout_extract.py "
        f"> '{wsl_path(log_path)}' 2>&1"
    )

    subprocess.run(
        ["wsl", "bash", "-lc", cmd],
        text=True
    )

    # -------------------------
    # ✅ OPEN ORIGINAL DESKTOP KLayout (NEW)
    # -------------------------
    try:
        klayout_exe = r"C:\Users\thefi\AppData\Roaming\KLayout\klayout_app.exe"
        subprocess.Popen([klayout_exe, gds_path])
    except Exception as e:
        print("KLayout desktop launch failed:", e)

    # -------------------------
    # Register artifacts
    # -------------------------
    register_klayout_artifacts(run)

    # -------------------------
    # Create presentation + slides
    # -------------------------
    presentation = auto_create_presentation(run)

    # -------------------------
    # Attach artifacts
    # -------------------------
    auto_attach_artifacts_to_slides(presentation, run)

    return JsonResponse({
        "ok": True,
        "message": "KLayout batch run completed",
        "run_id": str(run.id),
        "presentation_id": presentation.id,
        "redirect": f"/presentation/{presentation.id}/"
    })


def register_klayout_artifacts(run):
    base = Path(run.run_dir)  # ✅ already uploads/runs/<uuid>

    artifacts = [
        ("image", "Layout Preview", base / "preview.png"),
        ("metadata", "Layout Metadata", base / "metadata.json"),
        ("log", "KLayout Log", base / "klayout.log"),
    ]

    for kind, name, full_path in artifacts:
        if full_path.exists():
            RunArtifact.objects.create(
                run=run,
                artifact_type=(
                    "image" if kind == "image" else
                    "log" if kind == "log" else
                    "report"
                ),
                name=name,
                #  RELATIVE TO MEDIA_ROOT
                file_path=str(full_path.relative_to(settings.MEDIA_ROOT))
            )

from django.shortcuts import redirect, get_object_or_404
from .models import ToolRun, Presentation, Slide, SlideItem, RunArtifact

def create_presentation(request, run_id):
    run = get_object_or_404(ToolRun, id=run_id)

    #  Avoid duplicate presentations
    presentation = Presentation.objects.filter(run=run).first()
    if presentation:
        return redirect("presentation-detail", pk=presentation.id)

    # ✅ Create presentation
    presentation = Presentation.objects.create(
        title=f"KLayout Review – Run {run.id}",
        run=run,
        template=PresentationTemplate.objects.get(key="standard"),
        theme=PresentationTheme.objects.get(key="dark"),
    )

    # ✅ Slides
    layout_slide = Slide.objects.create(
        presentation=presentation,
        title="Layout View",
        order=1,
    )

    metadata_slide = Slide.objects.create(
        presentation=presentation,
        title="Metadata",
        order=2,
    )

    logs_slide = Slide.objects.create(
        presentation=presentation,
        title="Logs",
        order=3,
    )

    # ✅ Attach artifacts
    artifacts = RunArtifact.objects.filter(run=run)

    for artifact in artifacts:
        if artifact.artifact_type == "image":
            SlideItem.objects.create(
                slide=layout_slide,
                artifact=artifact,
                item_type="image",
            )

        elif artifact.artifact_type == "report":
            SlideItem.objects.create(
                slide=metadata_slide,
                artifact=artifact,
                item_type="attachment",
            )

        elif artifact.artifact_type == "log":
            SlideItem.objects.create(
                slide=logs_slide,
                artifact=artifact,
                item_type="log_snippet",
            )

    # ✅ Redirect to presentation UI
    return redirect("presentation-detail", pk=presentation.id)



from launcher.models import Presentation, Slide

from launcher.models import (
    Presentation,
    PresentationTemplate,
    PresentationTheme,
)
from django.db import transaction
from launcher.models import (
    Presentation,
    PresentationTemplate,
    PresentationTheme,
)

@transaction.atomic
def auto_create_presentation(run, user=None):
    """
    Automatically create a Presentation for a ToolRun.
    Ensures valid template/theme ForeignKey objects are used.
    """

    # Fetch default template
    default_template = PresentationTemplate.objects.filter(
        key="standard"
    ).first()

    # Fetch default theme
    default_theme = PresentationTheme.objects.filter(
        key="dark"
    ).first()

    # Hard safety checks (fail fast if DB is not seeded)
    if default_template is None:
        raise RuntimeError(
            "PresentationTemplate with key='standard' is missing. "
            "Seed templates before running tools."
        )

    if default_theme is None:
        raise RuntimeError(
            "PresentationTheme with key='dark' is missing. "
            "Seed themes before running tools."
        )

    # Create presentation
    presentation = Presentation.objects.create(
        title=f"KLayout Review – Run {run.id}",
        description="Auto-generated presentation from KLayout run",
        run=run,
        created_by=user,
        template=default_template,   # ✅ FK object
        theme=default_theme,         # ✅ FK object
    )

    return presentation



from launcher.models import Slide, SlideItem, RunArtifact

def auto_attach_artifacts_to_slides(presentation, run):
    from launcher.models import Slide, SlideItem, RunArtifact

    def get_slide(title, order):
        slide, _ = Slide.objects.get_or_create(
            presentation=presentation,
            title=title,
            defaults={"order": order},
        )
        return slide

    layout_slide = get_slide("Layout View", 1)
    metadata_slide = get_slide("Metadata", 2)
    logs_slide = get_slide("Logs", 3)

    artifacts = RunArtifact.objects.filter(run=run)

    for artifact in artifacts:
        atype = artifact.artifact_type.lower()

        # ---- IMAGE ----
        if atype == "image":
            SlideItem.objects.create(
                slide=layout_slide,
                item_type="image",
                artifact=artifact,
            )

        # ---- METADATA / REPORT ----
        elif atype == "report":
            SlideItem.objects.create(
                slide=metadata_slide,
                item_type="attachment",  # IMPORTANT
                artifact=artifact,
            )

        # ---- LOG ----
        elif atype == "log":
            SlideItem.objects.create(
                slide=logs_slide,
                item_type="log_snippet",
                artifact=artifact,
            )



from django.shortcuts import render, get_object_or_404
from .models import ToolRun

def run_detail(request, run_id):
    run = get_object_or_404(ToolRun, id=run_id)
    artifacts = run.artifacts.all()

    return render(
        request,
        "launcher/presentation/detail.html",
        {
            "run": run,
            "artifacts": artifacts,
            "is_run_view": True,
        }
    )


# =====================================================

def presentation_list(request):
    presentations = Presentation.objects.all().order_by("-created_at")

    return render(
        request,
        "launcher/presentation/list.html",
        {
            "presentations": presentations
        }
    )



# launcher/views.py
from pathlib import Path
from django.conf import settings
from django.shortcuts import get_object_or_404, render

from launcher.models import Presentation



from pathlib import Path

from django.conf import settings
from django.shortcuts import get_object_or_404, render

#from launcher.models import Presentation, TEMPLATE_CHOICES


from pathlib import Path

from django.conf import settings
from django.shortcuts import get_object_or_404, render

#from launcher.models import Presentation, TEMPLATE_CHOICES
from launcher.models import Presentation, PresentationTemplate, PresentationTheme


"""
def presentation_detail(request, pk):
    presentation = get_object_or_404(Presentation, pk=pk)

    # ===============================
    # HANDLE TEMPLATE / THEME UPDATE
    # ===============================
    if request.method == "POST":
        if "template_key" in request.POST:
            tpl = PresentationTemplate.objects.filter(
                key=request.POST["template_key"]
            ).first()
            if tpl:
                presentation.template = tpl
                presentation.save()

        if "theme_key" in request.POST:
            theme = PresentationTheme.objects.filter(
                key=request.POST["theme_key"]
            ).first()
            if theme:
                presentation.theme = theme
                presentation.save()

        return redirect(request.path)

    # ===============================
    # SLIDES
    # ===============================
    slides = presentation.slides.all().order_by("order")

    slide_id = request.GET.get("slide")
    selected_slide = slides.filter(id=slide_id).first() if slide_id else slides.first()

    # ===============================
    # ITEMS + INLINE CONTENT
    # ===============================
    items = []

    if selected_slide:
        for item in selected_slide.items.select_related("artifact"):
            item.inline_content = None

            if item.item_type in ("attachment", "log_snippet") and item.artifact:
                abs_path = Path(settings.MEDIA_ROOT) / item.artifact.file_path
                if abs_path.exists():
                    item.inline_content = abs_path.read_text(
                        encoding="utf-8", errors="replace"
                    )
                else:
                    item.inline_content = f"[FILE NOT FOUND] {abs_path}"

            items.append(item)

    # ===============================
    # TEMPLATE RESOLUTION (SAFE)
    # ===============================
    template_path = "launcher/presentation/detail.html"

    if presentation.template and presentation.template.base_template:
        template_path = presentation.template.base_template

    # ===============================
    # RENDER
    # ===============================
    return render(
        request,
        template_path,
        {
            "presentation": presentation,
            "slides": slides,
            "selected_slide": selected_slide,
            "items": items,

            # current selections
            "template": presentation.template,
            "theme": presentation.theme,

            # selectors
            "available_templates": PresentationTemplate.objects.all(),
            "available_themes": PresentationTheme.objects.all(),
        }
    )
"""
from pathlib import Path
from django.conf import settings
from django.shortcuts import get_object_or_404, render, redirect

from launcher.models import (
    Presentation,
    PresentationTemplate,
    PresentationTheme,
)


def presentation_detail(request, pk):
    presentation = get_object_or_404(Presentation, pk=pk)

    # ===== APPLY TEMPLATE / THEME =====
    if request.method == "POST":
        template_key = request.POST.get("template_key")
        theme_key = request.POST.get("theme_key")

        if template_key:
            tpl = PresentationTemplate.objects.filter(key=template_key).first()
            if tpl:
                presentation.template = tpl

        if theme_key:
            th = PresentationTheme.objects.filter(key=theme_key).first()
            if th:
                presentation.theme = th

        presentation.save()
        return redirect("presentation-detail", pk=presentation.pk)

    slides = presentation.slides.all().order_by("order")
    slide_id = request.GET.get("slide")
    selected_slide = slides.filter(id=slide_id).first() if slide_id else slides.first()

    items = []
    if selected_slide:
        for item in selected_slide.items.select_related("artifact").all():
            item.inline_content = None
            if item.item_type in ("attachment", "log_snippet") and item.artifact:
                abs_path = Path(settings.MEDIA_ROOT) / item.artifact.file_path
                if abs_path.exists():
                    item.inline_content = abs_path.read_text(
                        encoding="utf-8", errors="replace"
                    )
            items.append(item)

    return render(
        request,
        "launcher/presentation/detail.html",   # ✅ ALWAYS THIS
        {
            "presentation": presentation,
            "slides": slides,
            "selected_slide": selected_slide,
            "items": items,
            "available_templates": PresentationTemplate.objects.all(),
            "available_themes": PresentationTheme.objects.all(),
        },
    )



# launcher/views.py
from django.shortcuts import redirect, get_object_or_404

def set_presentation_template(request, pk):
    presentation = get_object_or_404(Presentation, pk=pk)

    if request.method == "POST":
        template = request.POST.get("template")
        if template:
            presentation.template = template
            presentation.save(update_fields=["template"])

    return redirect("presentation-detail", pk=pk)








from django.utils import timezone
from .models import Presentation


def presentation_create(request):
    presentation = Presentation.objects.create(
        title="New Design Review",
        created_by=request.user if request.user.is_authenticated else None,
        created_at=timezone.now()
    )

    return redirect("presentation-detail", pk=presentation.id)


from django.shortcuts import get_object_or_404
from django.http import HttpResponse
from django.template.loader import render_to_string
from django.utils import timezone
from pathlib import Path
from weasyprint import HTML

from django.conf import settings
from .models import Presentation



from pathlib import Path

def presentation_pdf(request, pk):
    presentation = get_object_or_404(Presentation, pk=pk)
    slides = presentation.slides.prefetch_related("items__artifact").all()

    for slide in slides:
        for item in slide.items.all():
            item.inline_content = ""

            if item.artifact:
                # ✅ Normalize Windows paths to URL paths
                item.artifact.url_path = item.artifact.file_path.replace("\\", "/")

            if item.item_type in ("attachment", "log_snippet") and item.artifact:
                abs_path = Path(settings.MEDIA_ROOT) / item.artifact.file_path
                if abs_path.exists():
                    item.inline_content = abs_path.read_text(
                        encoding="utf-8", errors="replace"
                    )

    html = render_to_string(
        "launcher/presentation/detail_pdf.html",
        {
            "presentation": presentation,
            "slides": slides,
            "request": request,
            "now": timezone.now(),
        }
    )

    pdf = HTML(
        string=html,
        base_url=settings.MEDIA_ROOT
    ).write_pdf()
    """
    return HttpResponse(
        pdf,
        content_type="application/pdf",
        #response["Content-Disposition"] = f'inline; filename="presentation-{pk}.pdf"'
        headers={
            "Content-Disposition": f'inline; filename="presentation-{pk}.pdf"'
        },
    )
    """
    response = HttpResponse(pdf, content_type="application/pdf")
    response["Content-Disposition"] = f'inline; filename="presentation-{pk}.pdf"'
    return response

# launcher/views.py

from pptx import Presentation as PPTPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from django.http import HttpResponse
from pathlib import Path
from django.conf import settings

def presentation_pptx(request, pk):
    presentation = get_object_or_404(Presentation, pk=pk)
    slides = presentation.slides.prefetch_related("items__artifact")

    prs = PPTPresentation()

    # Use 16:9 (standard PPT)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(11), Inches(2)
    )
    tf = title_box.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = presentation.title
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = f"Run {presentation.run.id} • {presentation.run.tool.name}"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # One REAL PPT slide per Slide model
    for slide in slides:
        ppt_slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Slide title
        title_box = ppt_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
        )
        title_tf = title_box.text_frame
        title_tf.text = slide.title
        title_tf.paragraphs[0].font.size = Pt(24)
        title_tf.paragraphs[0].font.bold = True

        y = 1.2  # vertical cursor

        for item in slide.items.all():

            # IMAGE → native PPT image
            if item.item_type == "image" and item.artifact:
                img_path = Path(settings.MEDIA_ROOT) / item.artifact.file_path
                if img_path.exists():
                    ppt_slide.shapes.add_picture(
                        str(img_path),
                        Inches(1),
                        Inches(y),
                        width=Inches(10)
                    )
                    y += 4.8

            # METADATA / LOGS → editable text box
            if item.item_type in ("attachment", "log_snippet") and item.artifact:
                path = Path(settings.MEDIA_ROOT) / item.artifact.file_path
                if path.exists():
                    content = path.read_text(errors="replace")

                    box = ppt_slide.shapes.add_textbox(
                        Inches(0.5), Inches(y), Inches(12), Inches(2)
                    )
                    tf = box.text_frame
                    tf.word_wrap = True
                    tf.text = content[:6000]  # PPT safety limit
                    tf.paragraphs[0].font.size = Pt(10)

                    y += 2.2

    # Return REAL PPTX
    response = HttpResponse(
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    response["Content-Disposition"] = (
        f'attachment; filename="presentation-{pk}.pptx"'
    )

    prs.save(response)
    return response





def logs_page(request):
    runs = ToolRun.objects.order_by("-created_at")[:50]
    return render(request, "launcher/logs.html", {"runs": runs})

def logs_home(request):
    """
    Landing page for logs.
    """
    recent_runs = ToolRun.objects.order_by("-created_at")[:20]

    return render(request, "launcher/logs_home.html", {
        "recent_runs": recent_runs
    })


def logs_by_scope(request, scope):
    """
    Logs filtered by scope:
    simulation | layout | runs | errors
    """

    qs = ToolRun.objects.order_by("-created_at")

    if scope == "simulation":
        qs = qs.filter(tool__slug="verilator")

    elif scope == "layout":
        qs = qs.filter(tool__slug="klayout")

    elif scope == "errors":
        qs = qs.filter(status="failed")

    elif scope == "runs":
        pass  # all runs

    else:
        return HttpResponse("Invalid log scope", status=404)

    return render(request, "launcher/logs_list.html", {
        "scope": scope,
        "runs": qs[:50]
    })

import subprocess

def run_bash(command, cwd=None, timeout=300):
    """
    Execute a shell command and return (stdout, stderr)
    """
    proc = subprocess.Popen(
        command,
        shell=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        cwd=cwd,
        text=True
    )

    try:
        stdout, stderr = proc.communicate(timeout=timeout)
    except subprocess.TimeoutExpired:
        proc.kill()
        stdout, stderr = proc.communicate()
        stderr += "\n[ERROR] Command timed out"

    return stdout, stderr

def windows_to_wsl(path):
    path = path.replace("\\", "/")
    if path[1:3] == ":/":
        drive = path[0].lower()
        return f"/mnt/{drive}{path[2:]}"
    return path



from django.utils import timezone

def execute_verilator_run(*, tool, user, upload_path):
    from django.utils import timezone
    import os

    run = ToolRun.objects.create(
        tool=tool,
        user=user,
        input_file=os.path.basename(upload_path),
        status="running"
    )

    try:
        wsl_file = windows_to_wsl(upload_path)

        cmd = f"verilator --lint-only {wsl_file}"

        stdout, stderr = run_bash(
            f"wsl bash -lc '{cmd}'"
        )

        run.stdout = stdout
        run.stderr = stderr
        run.status = "success" if not stderr else "failed"

    except Exception as e:
        run.stderr = str(e)
        run.status = "failed"

    run.completed_at = timezone.now()
    run.save()

    return run


@api_view(["POST"])
def run_tool(request, slug):
    tool = get_object_or_404(Tool, slug=slug)
    upload = request.FILES.get("file")

    if not upload:
        return Response({"ok": False, "error": "No file uploaded"}, status=400)

    uploads_dir = os.path.join(settings.BASE_DIR, "uploads")
    os.makedirs(uploads_dir, exist_ok=True)

    filename = f"{uuid.uuid4().hex}_{upload.name}"
    full_path = os.path.join(uploads_dir, filename)

    with open(full_path, "wb") as f:
        for chunk in upload.chunks():
            f.write(chunk)

    run = execute_verilator_run(
        tool=tool,
        user=request.user if request.user.is_authenticated else None,
        upload_path=full_path
    )

    return Response({
        "ok": run.status == "success",
        "run_id": run.id,
        "stdout": run.stdout,
        "stderr": run.stderr,
        "status": run.status
    })
def view_run_logs(request, run_id):
    run = get_object_or_404(ToolRun, id=run_id)

    return render(
        request,
        "launcher/run_detail.html",   # ✅ CORRECT PATH
        {"run": run}
    )
"""
def view_run_logs(request, run_id):
    run = get_object_or_404(
        ToolRun.objects.select_related("tool", "user"),
        id=run_id
    )

    return render(request, "launcher/logs/run_detail.html", {
        "run": run
    })
"""
def download_run_logs(request, run_id):
    run = get_object_or_404(ToolRun, id=run_id)

    content = f"""
========================================
EDA TOOL RUN LOG
========================================

Tool      : {run.tool.name}
Status    : {run.status}
Started   : {run.created_at}
Completed : {run.completed_at or "Running"}

----------------------------------------
STDOUT
----------------------------------------
{run.stdout or "No STDOUT"}

----------------------------------------
STDERR
----------------------------------------
{run.stderr or "No STDERR"}
"""

    response = HttpResponse(content, content_type="text/plain")
    response["Content-Disposition"] = (
        f'attachment; filename="{run.tool.slug}_run_{run.id}.log"'
    )
    return response

from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from django.shortcuts import get_object_or_404

from .models import Slide
from .serializers import SlideItemCreateSerializer


class AddSlideItemAPIView(APIView):
    def post(self, request, slide_id):
        slide = get_object_or_404(Slide, id=slide_id)

        # 🔒 Permission check (simple MVP)
        if slide.presentation.created_by != request.user:
            return Response(
                {"detail": "Not allowed"},
                status=status.HTTP_403_FORBIDDEN
            )

        serializer = SlideItemCreateSerializer(
            data=request.data,
            context={
                "slide": slide,
                "request": request,
            }
        )

        if serializer.is_valid():
            item = serializer.save()
            return Response(
                {
                    "status": "added",
                    "item_id": str(item.id)
                },
                status=status.HTTP_201_CREATED
            )

        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
